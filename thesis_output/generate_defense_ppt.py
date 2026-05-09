# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUT_DIR = Path(__file__).resolve().parent
OUT_PPT = OUT_DIR / "校园学术资源传承平台_答辩PPT.pptx"

WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)

BLUE = RGBColor(31, 78, 121)
LIGHT_BLUE = RGBColor(225, 237, 250)
MID_BLUE = RGBColor(68, 114, 196)
GREEN = RGBColor(91, 155, 120)
LIGHT_GREEN = RGBColor(232, 244, 235)
ORANGE = RGBColor(237, 125, 49)
LIGHT_ORANGE = RGBColor(255, 240, 220)
GRAY = RGBColor(88, 96, 105)
LIGHT_GRAY = RGBColor(246, 248, 250)
DARK = RGBColor(32, 38, 46)
WHITE = RGBColor(255, 255, 255)


def set_run(run, size=18, color=DARK, bold=False, font="微软雅黑"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, text, x, y, w, h, size=18, color=DARK, bold=False,
             align=PP_ALIGN.LEFT, fill=None, line=None, margin=0.08):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    p.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = p.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is not None:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_box(slide, text, x, y, w, h, fill=LIGHT_GRAY, line=MID_BLUE,
            size=16, bold=False, color=DARK, radius=True):
    typ = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(typ, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.25)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = p.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold)
    return shape


def add_line(slide, x1, y1, x2, y2, color=GRAY, width=1.4, arrow=False):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if arrow:
        conn.line.end_arrowhead = True
    return conn


def add_header(slide, title, section="校园学术资源传承平台"):
    add_text(slide, title, 0.55, 0.28, 8.7, 0.42, size=22, color=BLUE, bold=True)
    add_text(slide, section, 10.0, 0.32, 2.8, 0.34, size=10, color=GRAY, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.82), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(210, 218, 228)
    line.line.fill.background()


def add_footer(slide, idx):
    add_text(slide, f"{idx:02d}", 12.35, 7.08, 0.4, 0.22, size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def add_bullets(slide, items, x, y, w, h, size=17, color=DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        p.text = item
        p.font.name = "微软雅黑"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_table(slide, headers, rows, x, y, w, h, font_size=12):
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)
    )
    table = table_shape.table
    for col in range(len(headers)):
        table.columns[col].width = Inches(w / len(headers))
    for j, head in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        cell.text = head
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            set_run(r, size=font_size, color=WHITE, bold=True)
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT_GRAY
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT
            for r in p.runs:
                set_run(r, size=font_size, color=DARK)
    return table_shape


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def title_slide(prs):
    slide = blank(prs)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(245, 248, 252)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), WIDE_W, Inches(0.22)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = BLUE
    slide.shapes[-1].line.fill.background()
    add_text(slide, "校园学术资源传承平台的设计与实现", 1.0, 1.55, 11.3, 0.65, size=34, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "本科毕业论文（设计）答辩", 3.75, 2.45, 5.8, 0.45, size=22, color=DARK, align=PP_ALIGN.CENTER)
    add_box(slide, "二手书流转", 2.15, 3.45, 2.0, 0.62, fill=LIGHT_BLUE, line=MID_BLUE, size=17, bold=True)
    add_box(slide, "批注传承", 4.55, 3.45, 2.0, 0.62, fill=LIGHT_GREEN, line=GREEN, size=17, bold=True)
    add_box(slide, "学习路径", 6.95, 3.45, 2.0, 0.62, fill=LIGHT_ORANGE, line=ORANGE, size=17, bold=True)
    add_box(slide, "资源共享", 9.35, 3.45, 2.0, 0.62, fill=LIGHT_GRAY, line=GRAY, size=17, bold=True)
    add_text(slide, "学生：            指导教师：            专业：软件工程", 2.15, 5.55, 9.0, 0.35, size=15, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "2026 年 5 月", 5.2, 6.15, 2.9, 0.28, size=14, color=GRAY, align=PP_ALIGN.CENTER)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H

    title_slide(prs)

    # 2
    slide = blank(prs)
    add_header(slide, "答辩内容")
    topics = ["研究背景与意义", "需求分析", "系统设计", "系统实现", "系统测试", "总结与展望"]
    for i, t in enumerate(topics):
        x = 1.1 + (i % 3) * 3.95
        y = 1.55 + (i // 3) * 2.05
        add_box(slide, f"{i+1}\n{t}", x, y, 3.05, 1.15, fill=WHITE, line=MID_BLUE, size=18, bold=True)
    add_footer(slide, 2)

    # 3
    slide = blank(prs)
    add_header(slide, "研究背景：校园二手书与学习经验断层")
    add_bullets(slide, [
        "高校课程学习会产生大量教材、教辅与专业参考书，课程结束后部分书籍长期闲置。",
        "传统转让依赖微信群、QQ群和线下跳蚤市场，存在信息覆盖、分类检索困难、交易效率低等问题。",
        "书籍中的批注、复习重点和配套资料具有学习价值，但常随书籍闲置或转卖而散失。",
        "本课题希望在二手书流转基础上，进一步实现学习经验和学术资源的持续传承。"
    ], 0.9, 1.35, 6.1, 4.7, size=18)
    add_box(slide, "闲置教材", 8.0, 1.55, 2.1, 0.75, fill=LIGHT_BLUE, line=MID_BLUE, size=18, bold=True)
    add_box(slide, "交易沟通", 8.0, 2.75, 2.1, 0.75, fill=LIGHT_ORANGE, line=ORANGE, size=18, bold=True)
    add_box(slide, "批注资料", 8.0, 3.95, 2.1, 0.75, fill=LIGHT_GREEN, line=GREEN, size=18, bold=True)
    add_box(slide, "学习路径", 8.0, 5.15, 2.1, 0.75, fill=LIGHT_GRAY, line=GRAY, size=18, bold=True)
    add_line(slide, 10.1, 1.92, 11.25, 5.52, color=BLUE, width=2.0, arrow=True)
    add_box(slide, "学术资源传承平台", 10.75, 3.05, 1.75, 1.35, fill=WHITE, line=BLUE, size=15, bold=True)
    add_footer(slide, 3)

    # 4
    slide = blank(prs)
    add_header(slide, "研究现状与本文定位")
    rows = [
        ("校园二手交易", "关注物品发布、分类检索、沟通交易", "对学习批注与课程资源关注不足"),
        ("微信小程序平台", "入口轻量、无需安装、适合学生移动端使用", "需要结合具体校园业务设计功能闭环"),
        ("二手书交易", "降低学习成本，推动书籍循环利用", "多停留在书籍买卖层面"),
        ("本文工作", "交易、批注、资源、路径统一设计", "突出书籍流转中的学习经验沉淀")
    ]
    add_table(slide, ["方向", "已有研究特点", "本文改进重点"], rows, 0.75, 1.3, 11.85, 3.6, font_size=13)
    add_text(slide, "核心定位：不是普通二手交易平台，而是面向校园教材流转的“书籍 + 批注 + 资源 + 路径”综合平台。", 1.0, 5.45, 11.3, 0.55, size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER, fill=LIGHT_BLUE, line=MID_BLUE)
    add_footer(slide, 4)

    # 5
    slide = blank(prs)
    add_header(slide, "研究目标与主要工作")
    add_bullets(slide, [
        "完成面向用户和管理员两类角色的需求分析。",
        "设计前后端分离的总体架构、功能模块和核心业务流程。",
        "围绕用户、书籍、订单、批注、资源、学习路径等实体完成数据库设计。",
        "实现用户端微信小程序和后台管理端的主要功能。",
        "对登录认证、书籍交易、批注路径、后台管理等核心功能进行测试。"
    ], 0.9, 1.15, 6.2, 5.2, size=18)
    for i, label in enumerate(["需求分析", "总体设计", "功能实现", "系统测试"]):
        add_box(slide, label, 8.0, 1.35 + i * 1.15, 2.4, 0.62, fill=WHITE, line=BLUE, size=17, bold=True)
        if i < 3:
            add_line(slide, 9.2, 1.98 + i * 1.15, 9.2, 2.42 + i * 1.15, color=BLUE, arrow=True, width=1.8)
    add_footer(slide, 5)

    # 6
    slide = blank(prs)
    add_header(slide, "需求分析：用户与管理员")
    add_box(slide, "用户端", 0.9, 1.2, 5.4, 0.55, fill=LIGHT_BLUE, line=MID_BLUE, size=19, bold=True)
    add_box(slide, "管理端", 7.0, 1.2, 5.4, 0.55, fill=LIGHT_ORANGE, line=ORANGE, size=19, bold=True)
    add_bullets(slide, [
        "微信登录与个人资料维护",
        "学生认证、地址与通知管理",
        "书籍浏览、搜索、发布与收藏",
        "订单创建、模拟支付、发货、确认收货",
        "批注创建、学习路径、资源查看",
        "社区评论、举报与私信沟通"
    ], 1.05, 2.0, 5.0, 4.5, size=15)
    add_bullets(slide, [
        "后台登录与数据概览",
        "用户状态管理与认证审核",
        "书籍、资源、学习路径管理",
        "订单与纠纷辅助处理",
        "帖子、评论、举报与反馈处理",
        "轮播图和运营内容维护"
    ], 7.15, 2.0, 5.0, 4.5, size=15)
    add_footer(slide, 6)

    # 7
    slide = blank(prs)
    add_header(slide, "技术路线与项目结构")
    tech = [
        ("用户端", "UniApp / 微信小程序"),
        ("管理端", "Vue 3 / Web 后台"),
        ("后端", "Spring Boot / REST API"),
        ("持久层", "MyBatis-Flex / Mapper"),
        ("数据存储", "MySQL / Redis"),
        ("文件服务", "本地 uploads 静态资源")
    ]
    for i, (k, v) in enumerate(tech):
        x = 0.95 + (i % 2) * 5.9
        y = 1.25 + (i // 2) * 1.45
        add_box(slide, k, x, y, 1.55, 0.58, fill=BLUE, line=BLUE, size=15, bold=True, color=WHITE)
        add_box(slide, v, x + 1.7, y, 3.55, 0.58, fill=WHITE, line=RGBColor(205, 214, 224), size=15)
    add_text(slide, "选型原则：贴合项目实际、前后端分离、便于展示软件工程分层设计。", 1.0, 6.0, 11.3, 0.38, size=17, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 7)

    # 8 architecture
    slide = blank(prs)
    add_header(slide, "系统总体架构设计")
    layers = [
        ("表现层", ["微信小程序（UniApp）", "Web 管理后台（Vue 3）"], LIGHT_BLUE, MID_BLUE),
        ("业务层", ["用户认证", "书籍交易", "批注传承", "学习路径", "资源管理", "社区私信"], LIGHT_ORANGE, ORANGE),
        ("数据层", ["MySQL 8.0", "Redis 7.0"], LIGHT_GREEN, GREEN),
        ("支撑层", ["本地文件存储", "微信开放平台接口", "公共基础能力"], LIGHT_GRAY, GRAY),
    ]
    y = 1.15
    centers = []
    for idx, (name, items, fill, line) in enumerate(layers):
        add_box(slide, name, 0.75, y, 1.0, 0.78, fill=line, line=line, size=15, bold=True, color=WHITE)
        add_box(slide, "", 1.9, y, 10.75, 0.78, fill=fill, line=line, size=1)
        gap = 10.1 / len(items)
        for j, it in enumerate(items):
            add_box(slide, it, 2.1 + j * gap, y + 0.16, min(1.55, gap - 0.12), 0.45, fill=WHITE, line=line, size=11, bold=True)
        centers.append((7.25, y + 0.39))
        y += 1.38
    add_line(slide, 7.25, centers[0][1] + 0.42, 7.25, centers[1][1] - 0.42, color=BLUE, arrow=True, width=1.8)
    add_line(slide, 7.35, centers[1][1] - 0.42, 7.35, centers[0][1] + 0.42, color=BLUE, arrow=True, width=1.0)
    add_line(slide, 7.25, centers[1][1] + 0.42, 7.25, centers[2][1] - 0.42, color=BLUE, arrow=True, width=1.8)
    add_line(slide, 7.35, centers[2][1] - 0.42, 7.35, centers[1][1] + 0.42, color=BLUE, arrow=True, width=1.0)
    add_line(slide, 4.5, centers[1][1] + 0.55, 4.5, centers[3][1] - 0.42, color=GRAY, arrow=True, width=1.4)
    add_line(slide, 4.6, centers[3][1] - 0.42, 4.6, centers[1][1] + 0.55, color=GRAY, arrow=True, width=1.0)
    add_text(slide, "说明：请求方向由表现层进入业务层，业务层访问数据层和支撑能力；响应数据按相反方向返回。", 1.1, 6.7, 11.2, 0.3, size=12, color=GRAY, align=PP_ALIGN.CENTER)
    add_footer(slide, 8)

    # 9 modules
    slide = blank(prs)
    add_header(slide, "系统功能模块设计")
    add_box(slide, "校园学术资源传承平台", 4.55, 1.05, 4.1, 0.55, fill=WHITE, line=BLUE, size=19, bold=True)
    add_line(slide, 6.6, 1.6, 6.6, 2.12, color=BLUE, width=1.8)
    add_box(slide, "用户端（小程序）", 1.1, 2.12, 4.9, 0.52, fill=LIGHT_BLUE, line=MID_BLUE, size=17, bold=True)
    add_box(slide, "管理端（Web）", 7.3, 2.12, 4.9, 0.52, fill=LIGHT_ORANGE, line=ORANGE, size=17, bold=True)
    user_mods = ["首页浏览", "书籍模块", "交易模块", "传承模块", "社区模块", "个人中心"]
    admin_mods = ["仪表盘", "用户管理", "内容管理", "交易管理", "系统管理"]
    for i, m in enumerate(user_mods):
        add_box(slide, m, 1.2 + i * 0.78, 3.1, 0.62, 2.2, fill=WHITE, line=MID_BLUE, size=13, bold=True)
    for i, m in enumerate(admin_mods):
        add_box(slide, m, 7.55 + i * 0.9, 3.1, 0.72, 2.2, fill=WHITE, line=ORANGE, size=13, bold=True)
    add_footer(slide, 9)

    # 10 database
    slide = blank(prs)
    add_header(slide, "数据库设计：核心实体关系")
    entities = {
        "用户\nwx_user": (0.9, 1.35),
        "书籍\nbook": (4.0, 1.35),
        "订单\norder": (7.1, 1.35),
        "批注\nannotation": (4.0, 3.3),
        "资源\nresource": (7.1, 3.3),
        "学习路径\nlearning_path": (0.9, 3.3),
        "路径节点\npath_node": (0.9, 5.1),
        "社区/私信\npost/chat": (7.1, 5.1),
    }
    for name, (x, y) in entities.items():
        add_box(slide, name, x, y, 1.9, 0.76, fill=WHITE, line=BLUE, size=14, bold=True)
    add_line(slide, 2.8, 1.73, 4.0, 1.73, color=GRAY, arrow=True)
    add_line(slide, 5.9, 1.73, 7.1, 1.73, color=GRAY, arrow=True)
    add_line(slide, 4.95, 2.11, 4.95, 3.3, color=GRAY, arrow=True)
    add_line(slide, 5.9, 3.68, 7.1, 3.68, color=GRAY, arrow=True)
    add_line(slide, 1.85, 4.06, 1.85, 5.1, color=GRAY, arrow=True)
    add_line(slide, 2.8, 3.68, 4.0, 3.68, color=GRAY, arrow=True)
    add_line(slide, 8.05, 4.06, 8.05, 5.1, color=GRAY, arrow=True)
    add_text(slide, "关键表：wx_user、user_profile、book、order、annotation、resource、learning_path、path_node。", 1.1, 6.35, 11.1, 0.35, size=15, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 10)

    # 11 flows
    slide = blank(prs)
    add_header(slide, "核心业务流程")
    flows = [
        ("登录认证", ["微信授权", "生成 token", "访问受限功能"]),
        ("书籍交易", ["发布书籍", "买家沟通", "创建订单", "支付/发货/收货"]),
        ("批注传承", ["查看批注", "创建文字/图片批注", "点赞反馈"]),
        ("学习路径", ["创建路径", "维护节点", "查看资源", "记录进度"]),
    ]
    for r, (name, steps) in enumerate(flows):
        y = 1.22 + r * 1.28
        add_box(slide, name, 0.75, y, 1.45, 0.5, fill=BLUE, line=BLUE, size=14, bold=True, color=WHITE)
        for i, step in enumerate(steps):
            add_box(slide, step, 2.55 + i * 2.45, y, 1.55, 0.5, fill=WHITE, line=MID_BLUE, size=12)
            if i < len(steps) - 1:
                add_line(slide, 4.1 + i * 2.45, y + 0.25, 4.9 + i * 2.45, y + 0.25, color=GRAY, arrow=True)
    add_footer(slide, 11)

    # 12 implementation user
    slide = blank(prs)
    add_header(slide, "系统实现：用户端功能")
    add_bullets(slide, [
        "首页：展示轮播图、热门书籍、推荐学习路径和社区动态。",
        "发布：支持上传书籍图片，填写 ISBN、书名、作者、价格、成色等信息。",
        "详情：查看书籍信息、批注内容、相关资源，并发起沟通或交易。",
        "学习路径：查看路径节点、资源列表和学习进度，支持收藏与节点完成记录。",
        "社区与私信：支持帖子浏览、评论、举报、交易私信和订单快捷操作。"
    ], 0.9, 1.15, 6.4, 5.2, size=17)
    phone = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.35), Inches(1.05), Inches(2.45), Inches(5.2))
    phone.fill.solid()
    phone.fill.fore_color.rgb = RGBColor(248, 250, 252)
    phone.line.color.rgb = BLUE
    add_text(slide, "薪传", 8.75, 1.32, 1.65, 0.3, size=20, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, "热门书籍", 8.65, 2.0, 1.85, 0.45, fill=LIGHT_BLUE, line=MID_BLUE, size=12, bold=True)
    add_box(slide, "学习路径", 8.65, 2.7, 1.85, 0.45, fill=LIGHT_GREEN, line=GREEN, size=12, bold=True)
    add_box(slide, "社区动态", 8.65, 3.4, 1.85, 0.45, fill=LIGHT_ORANGE, line=ORANGE, size=12, bold=True)
    add_box(slide, "个人中心", 8.65, 4.1, 1.85, 0.45, fill=WHITE, line=GRAY, size=12, bold=True)
    add_box(slide, "发布书籍", 8.65, 4.8, 1.85, 0.45, fill=WHITE, line=GRAY, size=12, bold=True)
    add_footer(slide, 12)

    # 13 implementation admin
    slide = blank(prs)
    add_header(slide, "系统实现：管理端功能")
    rows = [
        ("后台登录", "管理员登录、菜单获取、退出登录"),
        ("数据概览", "统计用户、书籍、订单、内容等运行情况"),
        ("认证审核", "查看学生认证资料，处理审核结果"),
        ("内容管理", "书籍、资源、学习路径、帖子和评论状态维护"),
        ("运营处理", "举报、反馈、轮播图和订单纠纷辅助处理"),
    ]
    add_table(slide, ["模块", "主要功能"], rows, 1.0, 1.35, 11.3, 3.8, font_size=14)
    add_text(slide, "管理端重点承担平台治理和内容维护，不夸大为复杂运营系统，保持与项目实现一致。", 1.0, 5.75, 11.3, 0.42, size=17, color=BLUE, bold=True, align=PP_ALIGN.CENTER, fill=LIGHT_BLUE, line=MID_BLUE)
    add_footer(slide, 13)

    # 14 test
    slide = blank(prs)
    add_header(slide, "系统测试")
    rows = [
        ("用户登录与认证", "微信登录、受限页面访问、学生认证提交", "通过"),
        ("书籍发布与交易", "发布校验、订单创建、状态流转", "通过"),
        ("批注与学习路径", "批注创建、路径发布、节点进度更新", "通过"),
        ("后台管理", "管理员登录、认证审核、内容状态维护", "通过"),
    ]
    add_table(slide, ["测试模块", "测试内容", "结果"], rows, 0.9, 1.25, 11.5, 3.25, font_size=14)
    add_bullets(slide, [
        "测试方法：黑盒测试。",
        "测试目标：验证主要业务流程能否按预期完成。",
        "测试结论：系统核心功能完整，用户端与管理端主要操作均能正常反馈。"
    ], 1.1, 5.05, 10.8, 1.2, size=16, color=DARK)
    add_footer(slide, 14)

    # 15 summary
    slide = blank(prs)
    add_header(slide, "工作总结与特点")
    add_box(slide, "完成内容", 0.95, 1.25, 2.1, 0.55, fill=BLUE, line=BLUE, size=16, bold=True, color=WHITE)
    add_bullets(slide, [
        "完成需求分析、系统设计、数据库设计、功能实现和系统测试。",
        "实现用户端小程序与管理端后台的主要功能闭环。",
        "围绕书籍、批注、资源和学习路径建立核心数据模型。"
    ], 1.1, 2.0, 5.0, 2.5, size=16)
    add_box(slide, "系统特点", 7.1, 1.25, 2.1, 0.55, fill=GREEN, line=GREEN, size=16, bold=True, color=WHITE)
    add_bullets(slide, [
        "将二手书交易与学习经验传承结合。",
        "引入批注、资源、路径等学术资源管理能力。",
        "功能设计与实际项目控制器、页面和数据表保持一致。"
    ], 7.25, 2.0, 5.0, 2.5, size=16)
    add_footer(slide, 15)

    # 16 future thanks
    slide = blank(prs)
    add_header(slide, "不足与展望")
    add_bullets(slide, [
        "内容审核自动化程度仍可提升。",
        "批注图片可进一步接入 OCR 识别，提高知识提取效率。",
        "推荐功能可结合用户收藏、浏览和学习路径数据继续优化。",
        "后续可完善统一身份认证、消息通知和后台运营能力。"
    ], 1.05, 1.35, 5.9, 4.0, size=18)
    add_text(slide, "请各位老师批评指正", 7.25, 2.55, 4.8, 0.65, size=30, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "谢谢！", 8.65, 3.45, 2.0, 0.6, size=32, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 16)

    prs.save(OUT_PPT)
    return OUT_PPT


if __name__ == "__main__":
    print(build().resolve())
